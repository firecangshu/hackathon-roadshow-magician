# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 赛博朋克配色
BG_COLOR = RGBColor(13, 13, 13)  # #0D0D0D
NEON_PINK = RGBColor(255, 0, 255)  # #FF00FF
NEON_CYAN = RGBColor(0, 255, 255)  # #00FFFF

def add_slide(title, content_lines):
    """添加幻灯片"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = NEON_PINK
    p.alignment = PP_ALIGN.LEFT
    
    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(content_lines):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(36)
        p.font.color.rgb = NEON_CYAN
        p.space_after = Pt(20)
        p.alignment = PP_ALIGN.LEFT
    
    return slide

# Slide 1: 封面
add_slide("黑客松路演魔术师", [
    "把任意成果源代码一键转换为竞赛级路演材料",
    "",
    "7步工作流 · 19种产物 · 6维版本"
])

# Slide 2: 痛点
add_slide("99%的黑客松团队被路演拖后腿", [
    "DEMO做得好，但路演材料准备耗时过长",
    "浪费本该用于完善产品的宝贵时间",
    "",
    "解决方案：工程化生成，不浪费时间在排版和措辞上"
])

# Slide 3: 解决方案
add_slide("7步工作流 · 5个确认点", [
    "S1 源码深度分析 → S2 信息核对 → S3 路演定类",
    "→ S4 大纲构建 → S5 图表生成 → S6 成果生成 → S7 交付收尾",
    "",
    "每步有明确产出 + 确认点，避免方向跑偏返工"
])

# Slide 4: 产物体系
add_slide("六维版本 × 19种产物", [
    "六维：时长 × 侧重 × 风格 × 格式 × 语言 × 团队",
    "19种：赛前提报(4) + 现场展示(5) + 评委互动(4) + 赛后传播(4) + 存档复用(2)",
    "",
    "用户只需勾选，不浪费算力"
])

# Slide 5: 开源参与
add_slide("开源参与", [
    "MIT协议 · 社区友好 · 欢迎贡献",
    "",
    "Gitee地址：gitee.com/firecangshu/hackathon-roadshow-magician",
    "欢迎 Star · Issue · PR"
])

# 保存
output_path = r"C:\Users\User\.workbuddy\skills\hackathon-roadshow-magician__skillhub\路演材料\黑客松路演魔术师_赛博朋克风_5min.pptx"
prs.save(output_path)
print(f"PPT已生成：{output_path}")
